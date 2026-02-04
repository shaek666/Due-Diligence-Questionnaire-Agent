from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


@dataclass
class PageText:
    page_number: int
    text: str


def parse_pdf(path: str) -> list[PageText]:
    reader = PdfReader(path)
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append(PageText(page_number=index, text=text))
    return pages


def parse_docx(path: str) -> list[PageText]:
    doc = DocxDocument(path)
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    text = "\n".join(paragraphs)
    return [PageText(page_number=1, text=text)]


def parse_xlsx(path: str) -> list[PageText]:
    wb = load_workbook(path, data_only=True)
    pages: list[PageText] = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join([str(cell) for cell in row if cell is not None])
            if row_text:
                rows.append(row_text)
        pages.append(PageText(page_number=len(pages) + 1, text="\n".join(rows)))
    return pages


def parse_pptx(path: str) -> list[PageText]:
    presentation = Presentation(path)
    pages: list[PageText] = []
    for slide in presentation.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    texts.append(text)
        pages.append(PageText(page_number=len(pages) + 1, text="\n".join(texts)))
    return pages


def parse_document(path: str) -> list[PageText]:
    extension = Path(path).suffix.lower()
    if extension == ".pdf":
        return parse_pdf(path)
    if extension == ".docx":
        return parse_docx(path)
    if extension == ".xlsx":
        return parse_xlsx(path)
    if extension == ".pptx":
        return parse_pptx(path)
    raise ValueError(f"Unsupported file type: {extension}")
